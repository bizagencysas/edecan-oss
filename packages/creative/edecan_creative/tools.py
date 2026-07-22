"""Herramientas de creatividad: imágenes y documentos de oficina (`ARCHITECTURE.md`
§10.7; `ROADMAP_V2.md` §7.7, tabla `edecan_creative`).

Las cuatro herramientas (`generar_imagen`, `crear_documento`,
`crear_presentacion`, `crear_pdf`) generan un archivo en memoria y lo suben
con `edecan_creative._files.subir_archivo` (S3 + fila `files`), inyectado por
constructor con ese default real — el patrón inyectable que pide
`ROADMAP_V2.md` §7.7 para que los tests puedan sustituirlo sin tocar S3 ni
Postgres. Ninguna es `dangerous`: no publican nada, solo generan y guardan un
archivo privado del tenant (`ARCHITECTURE.md` §2, prefijo S3 por tenant).
Todas devuelven `data={"file_id", "filename"}` y un `content` en español listo
para mostrar en el chat.
"""

from __future__ import annotations

import io
import logging
import re
import unicodedata
from dataclasses import dataclass, field
from functools import lru_cache
from html.parser import HTMLParser
from typing import Any

from docx import Document
from edecan_core import Tool, ToolContext, ToolResult
from edecan_core.queue import enqueue
from fpdf import FPDF
from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation

from ._files import Uploader, subir_archivo
from .podcast import (
    FORMATOS_SOPORTADOS,
    MAX_SEGMENTOS,
    MIN_SEGMENTOS,
    GuionInvalidoError,
    generar_efecto,
    resolver_config_tts_tenant,
    validar_guion,
)
from .providers import DEFAULT_SIZE, ImageProvider, get_tenant_image_provider

logger = logging.getLogger(__name__)

# Límites para no generar archivos absurdos (ROADMAP_V2.md §7.7): topes sobre
# las listas de nivel superior (secciones/diapositivas/párrafos) y sobre las
# listas anidadas (párrafos por sección, bullets por diapositiva), más un tope
# de longitud por string individual. Ninguno de estos valores está pinned en
# el roadmap más allá de "máx 100 secciones/diapositivas, strings capped" — se
# eligieron generosos mientras acotados, documentados en `docs/creatividad.md`.
_MAX_TOP_ITEMS = 100
_MAX_NESTED_ITEMS = 200
_MAX_TITLE_CHARS = 300
_MAX_TEXT_CHARS = 4000
_MAX_PROMPT_PREVIEW_CHARS = 80

_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_PDF_MIME = "application/pdf"


def _cap_str(value: Any, max_chars: int) -> str:
    """Convierte `value` a `str`, recorta espacios y lo acota a `max_chars`."""
    return str(value or "").strip()[:max_chars]


def _slug(texto: str) -> str:
    """Nombre de archivo seguro derivado de `texto`: minúsculas, ASCII, `-` como
    separador (quita acentos con NFKD en vez de descartar la palabra completa,
    p. ej. "Métricas Q3" → "metricas-q3"). Nunca vacío: cae a "documento".
    """
    normalizado = unicodedata.normalize("NFKD", texto).encode("ascii", "ignore").decode("ascii")
    slug = re.sub(r"[^a-zA-Z0-9]+", "-", normalizado.strip().lower()).strip("-")
    return slug[:60] or "documento"


class GenerarImagenTool(Tool):
    name = "generar_imagen"
    description = (
        "Genera una imagen a partir de una descripción (prompt) usando el proveedor de "
        "imágenes configurado, y la guarda como archivo del usuario. El proveedor por "
        "defecto es un generador offline determinista (sin proveedor de imágenes real "
        "configurado, produce una imagen simple con el prompt como texto, no una foto "
        "realista)."
    )
    requires_flags = frozenset({"tools.images"})
    input_schema = {
        "type": "object",
        "properties": {
            "prompt": {
                "type": "string",
                "description": "Descripción de la imagen a generar.",
            },
            "tamano": {
                "type": "string",
                "description": "Tamaño de la imagen, formato ANCHOxALTO en píxeles.",
                "default": DEFAULT_SIZE,
            },
        },
        "required": ["prompt"],
    }

    def __init__(
        self,
        *,
        image_provider: ImageProvider | None = None,
        uploader: Uploader | None = None,
    ) -> None:
        self._image_provider = image_provider
        self._uploader = uploader or subir_archivo

    async def run(self, ctx: ToolContext, args: dict[str, Any]) -> ToolResult:
        prompt = _cap_str(args.get("prompt"), _MAX_TEXT_CHARS)
        if not prompt:
            return ToolResult(content="Necesito una descripción (prompt) para generar la imagen.")
        tamano = _cap_str(args.get("tamano"), 32) or DEFAULT_SIZE

        provider = self._image_provider or await get_tenant_image_provider(ctx)
        png_bytes = await provider.generate(prompt, size=tamano)

        filename = f"{_slug(prompt)}.png"
        file_id, filename = await self._uploader(
            ctx, data=png_bytes, filename=filename, mime="image/png"
        )

        preview = prompt[:_MAX_PROMPT_PREVIEW_CHARS]
        if len(prompt) > _MAX_PROMPT_PREVIEW_CHARS:
            preview += "…"
        return ToolResult(
            content=f"Creé la imagen «{preview}» ({tamano}) y la guardé como «{filename}».",
            data={
                "file_id": str(file_id),
                "filename": filename,
                "mime": "image/png",
                "alt_text": prompt,
                "caption": preview,
            },
        )


def _normalizar_secciones(raw: Any) -> list[dict[str, Any]]:
    if not isinstance(raw, list):
        return []
    secciones: list[dict[str, Any]] = []
    for item in raw[:_MAX_TOP_ITEMS]:
        if not isinstance(item, dict):
            continue
        encabezado = _cap_str(item.get("encabezado"), _MAX_TITLE_CHARS)
        if not encabezado:
            continue
        parrafos_raw = item.get("parrafos")
        parrafos = (
            [
                _cap_str(p, _MAX_TEXT_CHARS)
                for p in parrafos_raw[:_MAX_NESTED_ITEMS]
                if _cap_str(p, _MAX_TEXT_CHARS)
            ]
            if isinstance(parrafos_raw, list)
            else []
        )
        secciones.append({"encabezado": encabezado, "parrafos": parrafos})
    return secciones


class CrearDocumentoTool(Tool):
    name = "crear_documento"
    description = (
        "Genera un documento de Word (.docx) con un título y una o más secciones "
        "(encabezado + párrafos), y lo guarda como archivo del usuario."
    )
    input_schema = {
        "type": "object",
        "properties": {
            "titulo": {"type": "string", "description": "Título del documento."},
            "secciones": {
                "type": "array",
                "description": "Secciones del documento, cada una con encabezado y párrafos.",
                "items": {
                    "type": "object",
                    "properties": {
                        "encabezado": {"type": "string"},
                        "parrafos": {"type": "array", "items": {"type": "string"}},
                    },
                    "required": ["encabezado"],
                },
            },
        },
        "required": ["titulo", "secciones"],
    }

    def __init__(self, *, uploader: Uploader | None = None) -> None:
        self._uploader = uploader or subir_archivo

    async def run(self, ctx: ToolContext, args: dict[str, Any]) -> ToolResult:
        titulo = _cap_str(args.get("titulo"), _MAX_TITLE_CHARS)
        if not titulo:
            return ToolResult(content="Necesito un título para el documento.")
        if not isinstance(args.get("secciones"), list) or not args["secciones"]:
            return ToolResult(
                content=(
                    "Necesito al menos una sección (con encabezado y párrafos) para el documento."
                )
            )
        secciones = _normalizar_secciones(args["secciones"])
        if not secciones:
            return ToolResult(content="Ninguna de las secciones tenía un encabezado válido.")

        document = Document()
        document.add_heading(titulo, level=0)
        for seccion in secciones:
            document.add_heading(seccion["encabezado"], level=1)
            for parrafo in seccion["parrafos"]:
                document.add_paragraph(parrafo)

        buffer = io.BytesIO()
        document.save(buffer)
        data = buffer.getvalue()

        filename = f"{_slug(titulo)}.docx"
        file_id, filename = await self._uploader(ctx, data=data, filename=filename, mime=_DOCX_MIME)

        n_parrafos = sum(len(s["parrafos"]) for s in secciones)
        return ToolResult(
            content=(
                f"Creé el documento «{titulo}» con {len(secciones)} sección(es) y "
                f"{n_parrafos} párrafo(s), guardado como «{filename}»."
            ),
            data={"file_id": str(file_id), "filename": filename, "mime": _DOCX_MIME},
        )


def _normalizar_diapositivas(raw: Any) -> list[dict[str, Any]]:
    if not isinstance(raw, list):
        return []
    diapositivas: list[dict[str, Any]] = []
    for item in raw[:_MAX_TOP_ITEMS]:
        if not isinstance(item, dict):
            continue
        titulo = _cap_str(item.get("titulo"), _MAX_TITLE_CHARS)
        if not titulo:
            continue
        bullets_raw = item.get("bullets")
        bullets = (
            [
                _cap_str(b, _MAX_TEXT_CHARS)
                for b in bullets_raw[:_MAX_NESTED_ITEMS]
                if _cap_str(b, _MAX_TEXT_CHARS)
            ]
            if isinstance(bullets_raw, list)
            else []
        )
        diapositivas.append({"titulo": titulo, "bullets": bullets})
    return diapositivas


class CrearPresentacionTool(Tool):
    name = "crear_presentacion"
    description = (
        "Genera una presentación de PowerPoint (.pptx) con una portada de título y una "
        "diapositiva de título+contenido por cada elemento de 'diapositivas', y la "
        "guarda como archivo del usuario."
    )
    input_schema = {
        "type": "object",
        "properties": {
            "titulo": {"type": "string", "description": "Título de la presentación (portada)."},
            "diapositivas": {
                "type": "array",
                "description": "Diapositivas de contenido, cada una con título y viñetas.",
                "items": {
                    "type": "object",
                    "properties": {
                        "titulo": {"type": "string"},
                        "bullets": {"type": "array", "items": {"type": "string"}},
                    },
                    "required": ["titulo"],
                },
            },
        },
        "required": ["titulo", "diapositivas"],
    }

    def __init__(self, *, uploader: Uploader | None = None) -> None:
        self._uploader = uploader or subir_archivo

    async def run(self, ctx: ToolContext, args: dict[str, Any]) -> ToolResult:
        titulo = _cap_str(args.get("titulo"), _MAX_TITLE_CHARS)
        if not titulo:
            return ToolResult(content="Necesito un título para la presentación.")
        if not isinstance(args.get("diapositivas"), list) or not args["diapositivas"]:
            return ToolResult(
                content=(
                    "Necesito al menos una diapositiva (con título y viñetas) para la presentación."
                )
            )
        diapositivas = _normalizar_diapositivas(args["diapositivas"])
        if not diapositivas:
            return ToolResult(content="Ninguna de las diapositivas tenía un título válido.")

        presentation = Presentation()
        portada = presentation.slides.add_slide(presentation.slide_layouts[0])
        portada.shapes.title.text = titulo

        contenido_layout = presentation.slide_layouts[1]
        for diapositiva in diapositivas:
            slide = presentation.slides.add_slide(contenido_layout)
            slide.shapes.title.text = diapositiva["titulo"]
            bullets = diapositiva["bullets"] or [""]
            cuerpo = slide.placeholders[1].text_frame
            cuerpo.text = bullets[0]
            for bullet in bullets[1:]:
                cuerpo.add_paragraph().text = bullet

        buffer = io.BytesIO()
        presentation.save(buffer)
        data = buffer.getvalue()

        filename = f"{_slug(titulo)}.pptx"
        file_id, filename = await self._uploader(ctx, data=data, filename=filename, mime=_PPTX_MIME)

        return ToolResult(
            content=(
                f"Creé la presentación «{titulo}» con {len(diapositivas)} diapositiva(s) "
                f"de contenido (más la portada), guardada como «{filename}»."
            ),
            data={"file_id": str(file_id), "filename": filename, "mime": _PPTX_MIME},
        )


def _normalizar_parrafos_pdf(raw: Any) -> list[str]:
    if not isinstance(raw, list):
        return []
    return [
        _cap_str(p, _MAX_TEXT_CHARS) for p in raw[:_MAX_TOP_ITEMS] if _cap_str(p, _MAX_TEXT_CHARS)
    ]


def _sanitizar_pdf(texto: str) -> str:
    """Sanea `texto` al charset latin-1 (ISO-8859-1) que soportan las fuentes core
    de fpdf2 (Helvetica/Times/Courier), sustituyendo cualquier carácter fuera de
    rango por '?'. Evita descargar una fuente TrueType solo para soportar Unicode
    completo (regla dura de la plataforma: nada de red al generar un archivo) — la
    limitación de caracteres queda documentada en `docs/creatividad.md`.
    """
    # Las fuentes core de PDF no cubren toda la puntuacion Unicode. Convertimos
    # explicitamente los signos editoriales frecuentes antes de caer al
    # reemplazo final, para que un guion largo o unas comillas curvas nunca se
    # conviertan en un signo de interrogacion dentro del documento entregado.
    replacements = str.maketrans(
        {
            "—": "-",
            "–": "-",
            "“": '"',
            "”": '"',
            "‘": "'",
            "’": "'",
            "…": "...",
            "•": "-",
            "→": "->",
            "←": "<-",
        }
    )
    normalized = texto.translate(replacements)
    # Repara el residuo que dejaban versiones anteriores al convertir un guion
    # Unicode en `?` entre dos palabras. Los signos de interrogacion normales no
    # usan espacios a ambos lados en espanol.
    normalized = normalized.replace(" ? ", " - ")
    return normalized.encode("latin-1", errors="replace").decode("latin-1")


def _html_document_from(text: str) -> str | None:
    fenced = re.search(r"```(?:html)?\s*(<!doctype html.*?</html>)\s*```", text, re.I | re.S)
    if fenced:
        return fenced.group(1)
    direct = re.search(r"(<!doctype html.*?</html>)", text, re.I | re.S)
    return direct.group(1) if direct else None


@dataclass
class _DesignedPdfContent:
    title: str = ""
    subtitle: str = ""
    intro: str = ""
    chips: list[str] = field(default_factory=list)
    cards: list[tuple[str, str]] = field(default_factory=list)
    closing: str = ""


class _DesignedHtmlParser(HTMLParser):
    """Extrae la intención visual del HTML sin imprimir su código fuente."""

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.content = _DesignedPdfContent()
        self._skip_depth = 0
        self._capture_tag: str | None = None
        self._capture_class = ""
        self._capture_parts: list[str] = []
        self._pending_card_title: str | None = None

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        tag = tag.lower()
        if tag in {"style", "script", "head"}:
            self._skip_depth += 1
            return
        if self._skip_depth:
            return
        if tag == "br" and self._capture_tag:
            self._capture_parts.append("\n")
            return
        if self._capture_tag is None and tag in {"h1", "h2", "h3", "p", "span"}:
            self._capture_tag = tag
            self._capture_class = dict(attrs).get("class") or ""
            self._capture_parts = []

    def handle_endtag(self, tag: str) -> None:
        tag = tag.lower()
        if self._skip_depth:
            if tag in {"style", "script", "head"}:
                self._skip_depth -= 1
            return
        if self._capture_tag != tag:
            return
        text = " ".join("".join(self._capture_parts).split())
        css_class = self._capture_class.split()
        if text:
            if tag == "h1" and not self.content.title:
                self.content.title = text
            elif tag == "h3":
                self._pending_card_title = text
            elif tag == "span" and "chip" in css_class:
                self.content.chips.append(text)
            elif tag == "p" and "subtitle" in css_class:
                self.content.subtitle = text
            elif tag == "p" and "closing" in css_class:
                self.content.closing = text
            elif tag == "p" and self._pending_card_title:
                self.content.cards.append((self._pending_card_title, text))
                self._pending_card_title = None
            elif tag == "p" and not self.content.intro:
                self.content.intro = text
        self._capture_tag = None
        self._capture_class = ""
        self._capture_parts = []

    def handle_data(self, data: str) -> None:
        if not self._skip_depth and self._capture_tag:
            self._capture_parts.append(data)


def _parse_designed_html(source: str, fallback_title: str) -> _DesignedPdfContent:
    parser = _DesignedHtmlParser()
    parser.feed(source)
    parser.close()
    content = parser.content
    content.title = content.title or fallback_title
    return content


@lru_cache(maxsize=1)
def _aurora_background() -> bytes:
    width, height = 1240, 1754
    base = Image.new("RGBA", (width, height), (248, 248, 252, 255))
    glow = Image.new("RGBA", base.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(glow)
    draw.ellipse((-260, -300, 650, 630), fill=(133, 116, 255, 115))
    draw.ellipse((720, -180, 1440, 520), fill=(255, 164, 215, 100))
    draw.ellipse((650, 1080, 1500, 1920), fill=(99, 200, 255, 95))
    draw.ellipse((-360, 1180, 500, 1980), fill=(186, 223, 255, 105))
    glow = glow.filter(ImageFilter.GaussianBlur(125))
    composed = Image.alpha_composite(base, glow).convert("RGB")
    output = io.BytesIO()
    composed.save(output, format="PNG", optimize=True)
    return output.getvalue()


def _card(pdf: FPDF, x: float, y: float, width: float, height: float) -> None:
    pdf.set_fill_color(255, 255, 255)
    pdf.set_draw_color(231, 228, 244)
    pdf.rect(x, y, width, height, style="DF", round_corners=True)


def _render_designed_pdf(fallback_title: str, source: str) -> bytes:
    content = _parse_designed_html(source, fallback_title)
    pdf = FPDF(unit="mm", format="A4")
    pdf.set_auto_page_break(auto=False)
    pdf.add_page()
    pdf.image(io.BytesIO(_aurora_background()), x=0, y=0, w=210, h=297)

    pdf.set_text_color(108, 91, 239)
    pdf.set_font("Helvetica", "B", 8)
    pdf.set_xy(16, 18)
    pdf.cell(178, 5, "SISTEMA OPERATIVO COGNITIVO PERSONAL", align="C")

    pdf.set_text_color(32, 35, 51)
    pdf.set_font("Helvetica", "B", 31)
    pdf.set_xy(16, 28)
    pdf.cell(178, 15, _sanitizar_pdf(content.title)[:60], align="C")
    if content.subtitle:
        pdf.set_text_color(91, 92, 110)
        pdf.set_font("Helvetica", "I", 10)
        pdf.set_xy(24, 47)
        pdf.multi_cell(162, 5, _sanitizar_pdf(content.subtitle)[:260], align="C")

    intro = content.intro or "Una inteligencia personal que comprende, crea y ejecuta contigo."
    _card(pdf, 18, 65, 174, 45)
    pdf.set_text_color(68, 68, 86)
    pdf.set_font("Helvetica", size=9.5)
    pdf.set_xy(28, 74)
    pdf.multi_cell(154, 5.2, _sanitizar_pdf(intro)[:620], align="C")

    chips = content.chips[:10]
    if chips:
        pdf.set_text_color(111, 93, 210)
        pdf.set_font("Helvetica", "B", 7.5)
        rows = [chips[:5], chips[5:10]]
        y = 118.0
        for row in rows:
            if not row:
                continue
            widths = [
                min(32.0, max(18.0, pdf.get_string_width(_sanitizar_pdf(chip)) + 9))
                for chip in row
            ]
            gap = 3.0
            total = sum(widths) + gap * (len(widths) - 1)
            x = (210 - total) / 2
            for chip, width in zip(row, widths, strict=True):
                pdf.set_fill_color(246, 243, 255)
                pdf.set_draw_color(225, 218, 250)
                pdf.rect(x, y, width, 8, style="DF", round_corners=True)
                pdf.set_xy(x, y + 1.4)
                pdf.cell(width, 5, _sanitizar_pdf(chip)[:26], align="C")
                x += width + gap
            y += 10

    cards = content.cards[:6]
    if not cards:
        cards = [("Comprensión", intro), ("Ejecución", "Convierte intención en resultados útiles.")]
    start_y = 142 if len(chips) > 5 else 134
    card_width, card_height, gap_x, gap_y = 84.5, 34.0, 5.0, 5.0
    for index, (heading, body) in enumerate(cards):
        column, row = index % 2, index // 2
        is_unpaired_last = len(cards) % 2 == 1 and index == len(cards) - 1
        actual_width = 174.0 if is_unpaired_last else card_width
        x = 18 if is_unpaired_last else 18 + column * (card_width + gap_x)
        y = start_y + row * (card_height + gap_y)
        _card(pdf, x, y, actual_width, card_height)
        pdf.set_text_color(94, 75, 202)
        pdf.set_font("Helvetica", "B", 9)
        pdf.set_xy(x + 5, y + 5)
        pdf.cell(actual_width - 10, 5, _sanitizar_pdf(heading)[:70])
        pdf.set_text_color(82, 82, 99)
        pdf.set_font("Helvetica", size=7.6)
        pdf.set_xy(x + 5, y + 12)
        pdf.multi_cell(actual_width - 10, 4.2, _sanitizar_pdf(body)[:320])

    closing = content.closing or "Más claridad. Más capacidad. Mejores resultados."
    closing_y = min(269.0, start_y + ((len(cards) + 1) // 2) * (card_height + gap_y) + 5)
    pdf.set_text_color(85, 69, 185)
    pdf.set_font("Helvetica", "B", 11)
    pdf.set_xy(24, closing_y)
    pdf.multi_cell(162, 5.5, _sanitizar_pdf(closing)[:220], align="C")
    pdf.set_text_color(135, 128, 158)
    pdf.set_font("Helvetica", "B", 7)
    pdf.set_xy(16, 284)
    pdf.cell(178, 4, "EDECÁN", align="C")
    return bytes(pdf.output())


def _render_pdf(titulo: str, parrafos: list[str]) -> bytes:
    """Arma un PDF simple: portada con el título, luego el cuerpo en párrafos
    (fuente core Helvetica; `set_auto_page_break` agrega páginas automáticamente
    si el cuerpo no cabe en una sola).
    """
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)

    pdf.add_page()
    pdf.set_font("Helvetica", "B", 26)
    pdf.ln(70)
    pdf.multi_cell(0, 14, _sanitizar_pdf(titulo), align="C")

    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    for parrafo in parrafos:
        pdf.multi_cell(0, 7, _sanitizar_pdf(parrafo))
        pdf.ln(4)

    return bytes(pdf.output())


class CrearPdfTool(Tool):
    name = "crear_pdf"
    description = (
        "Genera un PDF con una portada (título) y un cuerpo de párrafos, y lo guarda "
        "como archivo del usuario. Cuando recibe HTML diseñado, extrae su intención "
        "visual y genera una página A4 pulida sin imprimir HTML o CSS como texto."
    )
    input_schema = {
        "type": "object",
        "properties": {
            "titulo": {"type": "string", "description": "Título / portada del PDF."},
            "parrafos": {
                "type": "array",
                "description": "Párrafos del cuerpo del documento, en orden.",
                "items": {"type": "string"},
            },
            "html": {
                "type": "string",
                "description": (
                    "HTML diseñado opcional; se renderiza visualmente, nunca como código."
                ),
            },
        },
        "required": ["titulo", "parrafos"],
    }

    def __init__(self, *, uploader: Uploader | None = None) -> None:
        self._uploader = uploader or subir_archivo

    async def run(self, ctx: ToolContext, args: dict[str, Any]) -> ToolResult:
        titulo = _cap_str(args.get("titulo"), _MAX_TITLE_CHARS)
        if not titulo:
            return ToolResult(content="Necesito un título para el PDF.")
        parrafos = _normalizar_parrafos_pdf(args.get("parrafos"))
        if not parrafos:
            return ToolResult(content="Necesito al menos un párrafo de contenido para el PDF.")

        html_source = _html_document_from(str(args.get("html") or ""))
        data = (
            _render_designed_pdf(titulo, html_source)
            if html_source
            else _render_pdf(titulo, parrafos)
        )

        filename = f"{_slug(titulo)}.pdf"
        file_id, filename = await self._uploader(ctx, data=data, filename=filename, mime=_PDF_MIME)

        return ToolResult(
            content=(
                f"Creé el PDF «{titulo}» con {len(parrafos)} párrafo(s), "
                f"guardado como «{filename}»."
            ),
            data={"file_id": str(file_id), "filename": filename, "mime": _PDF_MIME},
        )


# --- Podcasts y efectos de sonido (WP-V5-11, ARCHITECTURE.md §14) ---------------------------

_MAX_DESCRIPCION_EFECTO_CHARS = 500
_EFFECT_MIME = {"mp3": "audio/mpeg", "wav": "audio/wav"}


class CrearPodcastTool(Tool):
    name = "crear_podcast"
    description = (
        "Crea un podcast a partir de un guion (segmentos con orador + texto) y lo encola "
        "para generarse en segundo plano con el proveedor de voz (TTS) de tu cuenta "
        "(ElevenLabs si lo conectaste; si no, un audio de prueba offline) — aparece como "
        "archivo en unos minutos. El guion debe venir ya escrito en 'segmentos', esta "
        "herramienta no lo redacta."
    )
    requires_flags = frozenset({"tools.podcast"})
    input_schema = {
        "type": "object",
        "properties": {
            "titulo": {"type": "string", "description": "Título del podcast."},
            "segmentos": {
                "type": "array",
                "description": (
                    f"Guion del podcast: entre {MIN_SEGMENTOS} y {MAX_SEGMENTOS} segmentos, "
                    "cada uno con lo que dice un orador."
                ),
                "items": {
                    "type": "object",
                    "properties": {
                        "orador": {
                            "type": "string",
                            "description": "Nombre del orador de este segmento.",
                        },
                        "texto": {
                            "type": "string",
                            "description": "Lo que dice el orador en este segmento.",
                        },
                        "voice_id": {
                            "type": "string",
                            "description": (
                                "voice_id de ElevenLabs para este segmento (opcional; si se "
                                "omite se usa la voz por defecto de tu credencial)."
                            ),
                        },
                    },
                    "required": ["texto"],
                },
            },
            "formato": {
                "type": "string",
                "enum": list(FORMATOS_SOPORTADOS),
                "description": (
                    "Formato de salida preferido. El formato real siempre depende del "
                    "proveedor de voz que tengas conectado: 'wav' sin credencial (audio de "
                    "prueba), 'mp3' con ElevenLabs — pedir el que no coincide no fuerza una "
                    "conversión, se usa igual el formato real de tu proveedor."
                ),
            },
        },
        "required": ["titulo", "segmentos"],
    }

    async def run(self, ctx: ToolContext, args: dict[str, Any]) -> ToolResult:
        titulo = _cap_str(args.get("titulo"), _MAX_TITLE_CHARS)
        if not titulo:
            return ToolResult(content="Necesito un título para el podcast.")

        try:
            segmentos = validar_guion(args.get("segmentos"))
        except GuionInvalidoError as exc:
            return ToolResult(content=str(exc))

        formato = args.get("formato")
        if formato not in (None, *FORMATOS_SOPORTADOS):
            return ToolResult(
                content=f"'formato' debe ser {' o '.join(FORMATOS_SOPORTADOS)} (o vacío)."
            )

        payload = {
            "titulo": titulo,
            "formato": formato,
            "segmentos": [
                {"orador": s.orador, "texto": s.texto, "voice_id": s.voice_id} for s in segmentos
            ],
            # `JobEnvelope` no trae `user_id` (`ARCHITECTURE.md` §10.5) y este
            # job, a diferencia de `ingest_file`/`generate_content`, no tiene
            # ninguna fila previa (`file`/`conversation`) de la que heredarlo:
            # nace de cero e inserta una fila `files` nueva, cuya columna
            # `user_id` es `NOT NULL` (FK a `users`, migración `0001_initial`).
            # Viaja explícito en el payload para que el worker pueda atribuir
            # el archivo generado al usuario real que pidió el podcast.
            "user_id": str(ctx.user_id),
        }
        await enqueue(ctx.settings, "generate_podcast", payload, ctx.tenant_id)

        logger.info(
            "crear_podcast: encolado titulo=%r segmentos=%d tenant=%s",
            titulo,
            len(segmentos),
            ctx.tenant_id,
        )
        return ToolResult(
            content="Podcast en producción: aparecerá en Archivos en unos minutos.",
            data={"titulo": titulo, "segmentos": len(segmentos)},
        )


class GenerarEfectoSonidoTool(Tool):
    name = "generar_efecto_sonido"
    description = (
        "Genera un efecto de sonido SIN voz (ej. aplausos, lluvia, timbre, pasos) a partir "
        "de una descripción, usando el proveedor de voz (TTS) de tu cuenta (ElevenLabs si lo "
        "conectaste), y lo guarda como archivo del usuario. Sin credencial propia conectada "
        "genera un tono de prueba en vez de un efecto real."
    )
    requires_flags = frozenset({"tools.podcast"})
    input_schema = {
        "type": "object",
        "properties": {
            "descripcion": {
                "type": "string",
                "description": (
                    "Descripción del efecto de sonido a generar "
                    "(ej. 'aplausos entusiastas', 'lluvia suave sobre un techo')."
                ),
            },
        },
        "required": ["descripcion"],
    }

    def __init__(self, *, uploader: Uploader | None = None) -> None:
        self._uploader = uploader or subir_archivo

    async def run(self, ctx: ToolContext, args: dict[str, Any]) -> ToolResult:
        descripcion = _cap_str(args.get("descripcion"), _MAX_DESCRIPCION_EFECTO_CHARS)
        if not descripcion:
            return ToolResult(content="Necesito una descripción para generar el efecto de sonido.")

        cfg = await resolver_config_tts_tenant(
            session=getattr(ctx, "session", None),
            vault=getattr(ctx, "vault", None),
            tenant_id=getattr(ctx, "tenant_id", None),
        )
        audio = await generar_efecto(cfg, descripcion=descripcion, tenant_id=ctx.tenant_id)

        filename = f"{_slug(descripcion)}.{audio.formato}"
        file_id, filename = await self._uploader(
            ctx, data=audio.data, filename=filename, mime=_EFFECT_MIME[audio.formato]
        )

        aviso = (
            " (sin proveedor de voz conectado: generé un tono de prueba, no un efecto real)"
            if audio.es_stub
            else ""
        )
        return ToolResult(
            content=f"Generé el efecto «{descripcion}»{aviso} y lo guardé como «{filename}».",
            data={
                "file_id": str(file_id),
                "filename": filename,
                "mime": _EFFECT_MIME[audio.formato],
                "caption": descripcion,
                "es_stub": audio.es_stub,
            },
        )
