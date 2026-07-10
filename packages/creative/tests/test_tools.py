"""Tests de `edecan_creative.tools`: `generar_imagen`, `crear_documento`,
`crear_presentacion`, `crear_pdf` (`ROADMAP_V2.md` §7.7).

Cada tool recibe un `FakeUploader` inyectado por constructor (ver
`conftest.py`) — ninguno de estos tests toca S3 ni Postgres real. Los
archivos generados se re-abren con la misma librería que los produjo
(`python-docx`/`python-pptx`) o se valida su cabecera (`%PDF`, PNG) para
confirmar que la estructura es válida, no solo que hay bytes.
"""

from __future__ import annotations

import base64
import io
import json
from types import SimpleNamespace
from typing import Any
from uuid import uuid4

import edecan_creative.tools as tools_module
import httpx
import respx
from docx import Document
from edecan_creative.podcast import AudioGenerado
from edecan_creative.tools import (
    CrearDocumentoTool,
    CrearPdfTool,
    CrearPodcastTool,
    CrearPresentacionTool,
    GenerarEfectoSonidoTool,
    GenerarImagenTool,
)
from pptx import Presentation

PNG_SIGNATURE = b"\x89PNG\r\n\x1a\n"
_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


# --- generar_imagen ------------------------------------------------------------------


async def test_generar_imagen_usa_el_stub_por_defecto_y_sube_el_png(make_ctx, make_uploader):
    uploader = make_uploader()
    tool = GenerarImagenTool(uploader=uploader)

    resultado = await tool.run(
        make_ctx(), {"prompt": "un gato programador", "tamano": "256x256"}
    )

    assert len(uploader.llamadas) == 1
    llamada = uploader.llamadas[0]
    assert llamada["data"].startswith(PNG_SIGNATURE)
    assert llamada["mime"] == "image/png"
    assert llamada["filename"].endswith(".png")
    assert "gato programador" in resultado.content
    assert resultado.data == {"file_id": str(uploader.file_id), "filename": llamada["filename"]}


async def test_generar_imagen_usa_el_image_provider_inyectado(make_ctx, make_uploader):
    class _FakeProvider:
        def __init__(self) -> None:
            self.llamadas: list[tuple[str, str]] = []

        async def generate(self, prompt: str, size: str = "1024x1024") -> bytes:
            self.llamadas.append((prompt, size))
            return PNG_SIGNATURE + b"contenido-falso"

    provider = _FakeProvider()
    uploader = make_uploader()
    tool = GenerarImagenTool(image_provider=provider, uploader=uploader)

    await tool.run(make_ctx(), {"prompt": "un paisaje", "tamano": "512x512"})

    assert provider.llamadas == [("un paisaje", "512x512")]
    assert uploader.llamadas[0]["data"] == PNG_SIGNATURE + b"contenido-falso"


async def test_generar_imagen_sin_prompt_no_sube_nada(make_ctx, make_uploader):
    uploader = make_uploader()
    tool = GenerarImagenTool(uploader=uploader)

    resultado = await tool.run(make_ctx(), {"prompt": "   "})

    assert uploader.llamadas == []
    assert "descripción" in resultado.content.lower() or "prompt" in resultado.content.lower()


async def test_generar_imagen_prompt_largo_se_recorta_en_el_preview(make_ctx, make_uploader):
    uploader = make_uploader()
    tool = GenerarImagenTool(uploader=uploader)
    prompt_largo = "x" * 500

    resultado = await tool.run(make_ctx(), {"prompt": prompt_largo})

    assert "…" in resultado.content
    assert uploader.llamadas[0]["data"].startswith(PNG_SIGNATURE)


def test_generar_imagen_declara_flag_tools_images_y_no_es_dangerous():
    tool = GenerarImagenTool()
    assert tool.requires_flags == frozenset({"tools.images"})
    assert tool.dangerous is False


@respx.mock
async def test_generar_imagen_usa_la_credencial_bring_your_own_del_tenant(
    make_ctx, make_uploader, make_session, make_vault
):
    """Extremo a extremo (auditoría "riesgo-legal-tos"): si el tenant conectó
    su propia credencial de imágenes (`PUT /v1/credentials/images`), la tool
    la usa en vez de `IMAGES_API_KEY` de plataforma — sin `image_provider=`
    inyectado a mano, ejercitando `providers.get_tenant_image_provider` de
    verdad a través de `ctx.session`/`ctx.vault`."""
    fake_png = b"png-del-tenant"
    ruta = respx.post("https://images.tenant.example.com/v1/images/generations").mock(
        return_value=httpx.Response(
            200, json={"data": [{"b64_json": base64.b64encode(fake_png).decode("ascii")}]}
        )
    )

    bundle = SimpleNamespace(
        access_token=json.dumps(
            {
                "base_url": "https://images.tenant.example.com/v1",
                "api_key": "clave-del-tenant",
                "model": "modelo-del-tenant",
            }
        )
    )
    ctx = make_ctx(
        session=make_session([[{"id": "11111111-1111-1111-1111-111111111111"}]]),
        vault=make_vault(bundle=bundle),
    )
    uploader = make_uploader()
    tool = GenerarImagenTool(uploader=uploader)

    resultado = await tool.run(ctx, {"prompt": "un gato con la llave del tenant"})

    assert ruta.called
    assert uploader.llamadas[0]["data"] == fake_png
    assert "gato" in resultado.content


# --- crear_documento -------------------------------------------------------------------


async def test_crear_documento_genera_docx_valido_con_secciones_y_parrafos(make_ctx, make_uploader):
    uploader = make_uploader()
    tool = CrearDocumentoTool(uploader=uploader)

    resultado = await tool.run(
        make_ctx(),
        {
            "titulo": "Informe Trimestral",
            "secciones": [
                {"encabezado": "Resumen", "parrafos": ["Todo bien.", "Segundo párrafo."]},
                {"encabezado": "Riesgos", "parrafos": ["Ninguno crítico."]},
            ],
        },
    )

    assert len(uploader.llamadas) == 1
    llamada = uploader.llamadas[0]
    assert llamada["mime"] == _DOCX_MIME
    assert llamada["filename"].endswith(".docx")

    document = Document(io.BytesIO(llamada["data"]))
    textos = [p.text for p in document.paragraphs]
    assert "Informe Trimestral" in textos
    assert "Resumen" in textos
    assert "Todo bien." in textos
    assert "Segundo párrafo." in textos
    assert "Riesgos" in textos
    assert "Ninguno crítico." in textos

    assert "Informe Trimestral" in resultado.content
    assert "2 sección" in resultado.content
    assert resultado.data["file_id"] == str(uploader.file_id)


async def test_crear_documento_sin_titulo_no_sube_nada(make_ctx, make_uploader):
    uploader = make_uploader()
    tool = CrearDocumentoTool(uploader=uploader)

    resultado = await tool.run(make_ctx(), {"titulo": "  ", "secciones": [{"encabezado": "x"}]})

    assert uploader.llamadas == []
    assert "título" in resultado.content.lower()


async def test_crear_documento_sin_secciones_no_sube_nada(make_ctx, make_uploader):
    uploader = make_uploader()
    tool = CrearDocumentoTool(uploader=uploader)

    resultado = await tool.run(make_ctx(), {"titulo": "T", "secciones": []})

    assert uploader.llamadas == []
    assert "sección" in resultado.content.lower()


async def test_crear_documento_secciones_sin_encabezado_valido_no_sube_nada(
    make_ctx, make_uploader
):
    uploader = make_uploader()
    tool = CrearDocumentoTool(uploader=uploader)

    resultado = await tool.run(make_ctx(), {"titulo": "T", "secciones": [{"encabezado": "   "}]})

    assert uploader.llamadas == []
    assert resultado.content


async def test_crear_documento_seccion_sin_parrafos_no_falla(make_ctx, make_uploader):
    uploader = make_uploader()
    tool = CrearDocumentoTool(uploader=uploader)

    resultado = await tool.run(
        make_ctx(), {"titulo": "T", "secciones": [{"encabezado": "Solo encabezado"}]}
    )

    assert len(uploader.llamadas) == 1
    assert "0 párrafo" in resultado.content


async def test_crear_documento_se_acota_a_100_secciones(make_ctx, make_uploader):
    uploader = make_uploader()
    tool = CrearDocumentoTool(uploader=uploader)
    secciones = [{"encabezado": f"S{i}", "parrafos": ["x"]} for i in range(150)]

    resultado = await tool.run(make_ctx(), {"titulo": "T", "secciones": secciones})

    assert "100 sección" in resultado.content


# --- crear_presentacion ---------------------------------------------------------------


async def test_crear_presentacion_genera_pptx_valido_con_portada_y_diapositivas(
    make_ctx, make_uploader
):
    uploader = make_uploader()
    tool = CrearPresentacionTool(uploader=uploader)

    resultado = await tool.run(
        make_ctx(),
        {
            "titulo": "Estrategia 2027",
            "diapositivas": [
                {"titulo": "Objetivos", "bullets": ["Crecer", "Retener"]},
                {"titulo": "Riesgos", "bullets": ["Competencia"]},
            ],
        },
    )

    assert len(uploader.llamadas) == 1
    llamada = uploader.llamadas[0]
    assert llamada["mime"] == _PPTX_MIME
    assert llamada["filename"].endswith(".pptx")

    presentation = Presentation(io.BytesIO(llamada["data"]))
    assert len(presentation.slides) == 3  # portada + 2 diapositivas de contenido
    assert presentation.slides[0].shapes.title.text == "Estrategia 2027"
    assert presentation.slides[1].shapes.title.text == "Objetivos"
    bullets_slide_1 = [p.text for p in presentation.slides[1].placeholders[1].text_frame.paragraphs]
    assert bullets_slide_1 == ["Crecer", "Retener"]
    assert presentation.slides[2].shapes.title.text == "Riesgos"

    assert "Estrategia 2027" in resultado.content
    assert "2 diapositiva" in resultado.content


async def test_crear_presentacion_sin_titulo_no_sube_nada(make_ctx, make_uploader):
    uploader = make_uploader()
    tool = CrearPresentacionTool(uploader=uploader)

    resultado = await tool.run(
        make_ctx(), {"titulo": "", "diapositivas": [{"titulo": "x"}]}
    )

    assert uploader.llamadas == []
    assert "título" in resultado.content.lower()


async def test_crear_presentacion_sin_diapositivas_no_sube_nada(make_ctx, make_uploader):
    uploader = make_uploader()
    tool = CrearPresentacionTool(uploader=uploader)

    resultado = await tool.run(make_ctx(), {"titulo": "T", "diapositivas": []})

    assert uploader.llamadas == []
    assert "diapositiva" in resultado.content.lower()


async def test_crear_presentacion_se_acota_a_100_diapositivas(make_ctx, make_uploader):
    uploader = make_uploader()
    tool = CrearPresentacionTool(uploader=uploader)
    diapositivas = [{"titulo": f"D{i}", "bullets": ["x"]} for i in range(150)]

    resultado = await tool.run(make_ctx(), {"titulo": "T", "diapositivas": diapositivas})

    assert "100 diapositiva" in resultado.content


# --- crear_pdf -------------------------------------------------------------------------


async def test_crear_pdf_genera_pdf_valido(make_ctx, make_uploader):
    uploader = make_uploader()
    tool = CrearPdfTool(uploader=uploader)

    resultado = await tool.run(
        make_ctx(), {"titulo": "Reporte", "parrafos": ["Primer párrafo.", "Segundo párrafo."]}
    )

    assert len(uploader.llamadas) == 1
    llamada = uploader.llamadas[0]
    assert llamada["data"][:4] == b"%PDF"
    assert llamada["mime"] == "application/pdf"
    assert llamada["filename"].endswith(".pdf")
    assert "Reporte" in resultado.content
    assert "2 párrafo" in resultado.content


async def test_crear_pdf_sanea_caracteres_fuera_de_latin1_sin_reventar(make_ctx, make_uploader):
    uploader = make_uploader()
    tool = CrearPdfTool(uploader=uploader)

    resultado = await tool.run(
        make_ctx(),
        {"titulo": "Título con eñe y emoji 🚀", "parrafos": ["Raya larga — y comillas “curvas”."]},
    )

    assert len(uploader.llamadas) == 1
    assert uploader.llamadas[0]["data"][:4] == b"%PDF"
    assert resultado.data["file_id"] == str(uploader.file_id)


async def test_crear_pdf_sin_titulo_no_sube_nada(make_ctx, make_uploader):
    uploader = make_uploader()
    tool = CrearPdfTool(uploader=uploader)

    resultado = await tool.run(make_ctx(), {"titulo": "  ", "parrafos": ["x"]})

    assert uploader.llamadas == []
    assert "título" in resultado.content.lower()


async def test_crear_pdf_sin_parrafos_no_sube_nada(make_ctx, make_uploader):
    uploader = make_uploader()
    tool = CrearPdfTool(uploader=uploader)

    resultado = await tool.run(make_ctx(), {"titulo": "T", "parrafos": []})

    assert uploader.llamadas == []
    assert "párrafo" in resultado.content.lower()


# --- crear_podcast (WP-V5-11) ------------------------------------------------------------


def _install_fake_enqueue(monkeypatch) -> list[tuple[Any, str, dict, Any]]:
    """Mismo patrón que `packages/agents/tests/test_tools.py`: `enqueue` se
    monkeypatchea sobre el NOMBRE importado en `edecan_creative.tools`, así
    el test nunca abre una conexión SQS real."""
    llamadas: list[tuple[Any, str, dict, Any]] = []

    async def fake_enqueue(settings, job_type, payload, tenant_id, **kwargs):
        llamadas.append((settings, job_type, payload, tenant_id))
        return uuid4()

    monkeypatch.setattr(tools_module, "enqueue", fake_enqueue)
    return llamadas


def test_crear_podcast_declara_flag_tools_podcast_y_no_es_dangerous():
    tool = CrearPodcastTool()
    assert tool.requires_flags == frozenset({"tools.podcast"})
    assert tool.dangerous is False
    assert tool.input_schema["required"] == ["titulo", "segmentos"]


async def test_crear_podcast_sin_titulo_no_encola(make_ctx, monkeypatch):
    llamadas = _install_fake_enqueue(monkeypatch)

    resultado = await CrearPodcastTool().run(
        make_ctx(), {"titulo": "   ", "segmentos": [{"texto": "hola"}]}
    )

    assert "título" in resultado.content.lower()
    assert llamadas == []


async def test_crear_podcast_guion_invalido_no_encola(make_ctx, monkeypatch):
    llamadas = _install_fake_enqueue(monkeypatch)

    resultado = await CrearPodcastTool().run(make_ctx(), {"titulo": "T", "segmentos": []})

    assert resultado.content  # el mensaje viene de podcast.validar_guion
    assert llamadas == []


async def test_crear_podcast_segmento_sin_texto_no_encola(make_ctx, monkeypatch):
    llamadas = _install_fake_enqueue(monkeypatch)

    resultado = await CrearPodcastTool().run(
        make_ctx(), {"titulo": "T", "segmentos": [{"orador": "Ana"}]}
    )

    assert "texto" in resultado.content.lower()
    assert llamadas == []


async def test_crear_podcast_formato_invalido_no_encola(make_ctx, monkeypatch):
    llamadas = _install_fake_enqueue(monkeypatch)

    resultado = await CrearPodcastTool().run(
        make_ctx(), {"titulo": "T", "segmentos": [{"texto": "hola"}], "formato": "ogg"}
    )

    assert "formato" in resultado.content.lower()
    assert llamadas == []


async def test_crear_podcast_camino_feliz_encola_generate_podcast(make_ctx, monkeypatch):
    llamadas = _install_fake_enqueue(monkeypatch)
    tenant_id = uuid4()
    user_id = uuid4()
    ctx = make_ctx(tenant_id=tenant_id, user_id=user_id)

    resultado = await CrearPodcastTool().run(
        ctx,
        {
            "titulo": "Mi Podcast",
            "segmentos": [
                {"orador": "Ana", "texto": "Hola a todos", "voice_id": "voz-1"},
                {"texto": "Segundo segmento, sin orador declarado"},
            ],
            "formato": "mp3",
        },
    )

    assert "producción" in resultado.content.lower()
    assert resultado.data == {"titulo": "Mi Podcast", "segmentos": 2}

    assert len(llamadas) == 1
    _settings, job_type, payload, enq_tenant_id = llamadas[0]
    assert job_type == "generate_podcast"
    assert enq_tenant_id is tenant_id
    assert payload["titulo"] == "Mi Podcast"
    assert payload["formato"] == "mp3"
    assert payload["user_id"] == str(user_id)
    assert payload["segmentos"] == [
        {"orador": "Ana", "texto": "Hola a todos", "voice_id": "voz-1"},
        {"orador": "Orador 2", "texto": "Segundo segmento, sin orador declarado", "voice_id": None},
    ]


async def test_crear_podcast_formato_omitido_viaja_como_none_en_el_payload(make_ctx, monkeypatch):
    llamadas = _install_fake_enqueue(monkeypatch)

    await CrearPodcastTool().run(make_ctx(), {"titulo": "T", "segmentos": [{"texto": "hola"}]})

    _settings, _job_type, payload, _tenant_id = llamadas[0]
    assert payload["formato"] is None


async def test_crear_podcast_se_acota_a_30_segmentos_y_reporta_el_error(make_ctx, monkeypatch):
    llamadas = _install_fake_enqueue(monkeypatch)
    segmentos = [{"texto": f"segmento {i}"} for i in range(31)]

    resultado = await CrearPodcastTool().run(
        make_ctx(), {"titulo": "T", "segmentos": segmentos}
    )

    assert "30" in resultado.content
    assert llamadas == []


# --- generar_efecto_sonido (WP-V5-11) -----------------------------------------------------


def _install_fake_generar_efecto(monkeypatch, audio: AudioGenerado):
    async def fake_generar_efecto(cfg, *, descripcion, tenant_id=None):
        return audio

    monkeypatch.setattr(tools_module, "generar_efecto", fake_generar_efecto)


def test_generar_efecto_sonido_declara_flag_tools_podcast_y_no_es_dangerous():
    tool = GenerarEfectoSonidoTool()
    assert tool.requires_flags == frozenset({"tools.podcast"})
    assert tool.dangerous is False


async def test_generar_efecto_sonido_sin_descripcion_no_sube_nada(make_ctx, make_uploader):
    uploader = make_uploader()
    tool = GenerarEfectoSonidoTool(uploader=uploader)

    resultado = await tool.run(make_ctx(), {"descripcion": "   "})

    assert uploader.llamadas == []
    assert "descripción" in resultado.content.lower()


async def test_generar_efecto_sonido_stub_sube_wav_y_lo_avisa(make_ctx, make_uploader, monkeypatch):
    _install_fake_generar_efecto(
        monkeypatch, AudioGenerado(data=b"RIFF-wav-fake", formato="wav", es_stub=True)
    )
    uploader = make_uploader()
    tool = GenerarEfectoSonidoTool(uploader=uploader)

    resultado = await tool.run(make_ctx(), {"descripcion": "aplausos entusiastas"})

    assert len(uploader.llamadas) == 1
    llamada = uploader.llamadas[0]
    assert llamada["data"] == b"RIFF-wav-fake"
    assert llamada["mime"] == "audio/wav"
    assert llamada["filename"].endswith(".wav")
    assert "tono de prueba" in resultado.content
    assert resultado.data == {
        "file_id": str(uploader.file_id),
        "filename": llamada["filename"],
        "es_stub": True,
    }


async def test_generar_efecto_sonido_elevenlabs_sube_mp3_sin_avisar_stub(
    make_ctx, make_uploader, monkeypatch
):
    _install_fake_generar_efecto(
        monkeypatch, AudioGenerado(data=b"mp3-real-bytes", formato="mp3", es_stub=False)
    )
    uploader = make_uploader()
    tool = GenerarEfectoSonidoTool(uploader=uploader)

    resultado = await tool.run(make_ctx(), {"descripcion": "lluvia suave sobre un techo"})

    llamada = uploader.llamadas[0]
    assert llamada["mime"] == "audio/mpeg"
    assert llamada["filename"].endswith(".mp3")
    assert "tono de prueba" not in resultado.content
    assert resultado.data["es_stub"] is False


async def test_generar_efecto_sonido_resuelve_config_del_tenant_con_ctx(
    make_ctx, make_session, make_vault, make_uploader, monkeypatch
):
    llamadas_resolver: list[tuple[Any, Any, Any]] = []

    async def fake_resolver(*, session, vault, tenant_id):
        llamadas_resolver.append((session, vault, tenant_id))
        return None

    monkeypatch.setattr(tools_module, "resolver_config_tts_tenant", fake_resolver)
    _install_fake_generar_efecto(
        monkeypatch, AudioGenerado(data=b"x", formato="wav", es_stub=True)
    )

    session = make_session()
    vault = make_vault()
    tenant_id = uuid4()
    ctx = make_ctx(session=session, vault=vault, tenant_id=tenant_id)

    await GenerarEfectoSonidoTool(uploader=make_uploader()).run(ctx, {"descripcion": "algo"})

    assert llamadas_resolver == [(session, vault, tenant_id)]
