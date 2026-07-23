"""Lectura universal y edición reversible de archivos desde el chat."""

from __future__ import annotations

import io
import zipfile
from pathlib import PurePath
from typing import Any

from edecan_core import Tool, ToolContext, ToolResult

from . import _s3
from ._util import parse_uuid

_MAX_FILE_BYTES = 25 * 1024 * 1024
_MAX_TEXT_CHARS = 80_000
_MAX_ZIP_UNCOMPRESSED_BYTES = 100 * 1024 * 1024


def _normalized_mime(mime: str) -> str:
    return (mime or "application/octet-stream").split(";", 1)[0].strip().lower()


def _validate_office_archive(data: bytes) -> None:
    """Rechaza ZIP bombs antes de entregar OOXML a librerías de terceros."""

    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        expanded = sum(info.file_size for info in archive.infolist())
    if expanded > _MAX_ZIP_UNCOMPRESSED_BYTES:
        raise ValueError("El documento expandido supera 100 MB")


def _extract_text(data: bytes, *, filename: str, mime: str) -> str | None:
    name = filename.lower()
    normalized = _normalized_mime(mime)
    if normalized.startswith("text/") or name.endswith(
        (".txt", ".md", ".csv", ".tsv", ".json", ".xml", ".html", ".css", ".js", ".py")
    ):
        return data.decode("utf-8", errors="replace")
    if normalized == "application/pdf" or name.endswith(".pdf"):
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        pages: list[str] = []
        total_chars = 0
        for index, page in enumerate(reader.pages[:200], start=1):
            text = f"[Página {index}]\n{page.extract_text() or ''}"
            pages.append(text)
            total_chars += len(text)
            if total_chars >= _MAX_TEXT_CHARS:
                break
        return "\n\n".join(pages)
    if name.endswith(".docx") or normalized.endswith("wordprocessingml.document"):
        import docx

        _validate_office_archive(data)
        document = docx.Document(io.BytesIO(data))
        paragraphs: list[str] = []
        total_chars = 0
        for paragraph in document.paragraphs[:20_000]:
            paragraphs.append(paragraph.text)
            total_chars += len(paragraph.text)
            if total_chars >= _MAX_TEXT_CHARS:
                break
        return "\n".join(paragraphs)
    if name.endswith(".xlsx") or normalized.endswith("spreadsheetml.sheet"):
        from openpyxl import load_workbook

        _validate_office_archive(data)
        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        sheets: list[str] = []
        total_chars = 0
        for sheet in workbook.worksheets[:50]:
            rows: list[str] = []
            for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                rendered = "\t".join(
                    "" if value is None else str(value) for value in row[:200]
                )
                rows.append(rendered)
                total_chars += len(rendered)
                if row_index >= 10_000 or total_chars >= _MAX_TEXT_CHARS:
                    break
            sheets.append(f"[Hoja: {sheet.title}]\n" + "\n".join(rows))
            if total_chars >= _MAX_TEXT_CHARS:
                break
        return "\n\n".join(sheets)
    if name.endswith(".pptx") or normalized.endswith("presentationml.presentation"):
        from pptx import Presentation

        _validate_office_archive(data)
        presentation = Presentation(io.BytesIO(data))
        slides: list[str] = []
        total_chars = 0
        for index, slide in enumerate(presentation.slides, start=1):
            if index > 500:
                break
            text = "\n".join(
                shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)
            )
            slides.append(f"[Diapositiva {index}]\n{text}")
            total_chars += len(text)
            if total_chars >= _MAX_TEXT_CHARS:
                break
        return "\n\n".join(slides)
    return None


class LeerArchivoTool(Tool):
    name = "leer_archivo"
    description = (
        "Abre y lee un archivo privado ya adjunto: PDF, Word, PowerPoint, Excel, CSV, "
        "JSON, Markdown, código o texto. Si es una imagen, la analiza visualmente. "
        "Úsala antes de resumir, opinar, corregir o transformar un adjunto."
    )
    input_schema = {
        "type": "object",
        "properties": {
            "file_id": {"type": "string", "description": "Identificador del adjunto."},
            "pregunta": {
                "type": "string",
                "description": "Qué debe buscar o responder sobre el archivo.",
            },
        },
        "required": ["file_id"],
    }

    async def run(self, ctx: ToolContext, args: dict[str, Any]) -> ToolResult:
        file_id = parse_uuid(args.get("file_id"))
        if file_id is None:
            return ToolResult(content="'file_id' no es un identificador válido.")
        archivo = await _s3.descargar_archivo(ctx, file_id)
        if archivo is None:
            return ToolResult(content="No encontré ese archivo.")
        if archivo.size_bytes > _MAX_FILE_BYTES:
            return ToolResult(content="El archivo supera 25 MB y no puedo abrirlo completo.")
        if _normalized_mime(archivo.mime).startswith("image/"):
            from .vision import AnalizarImagenTool

            return await AnalizarImagenTool().run(
                ctx,
                {
                    "file_id": str(file_id),
                    "pregunta": str(args.get("pregunta") or "").strip()
                    or "Describe, transcribe y analiza esta imagen.",
                },
            )
        try:
            extracted = _extract_text(
                archivo.contenido, filename=archivo.filename, mime=archivo.mime
            )
        except Exception as exc:  # noqa: BLE001 - formatos de terceros heterogéneos
            return ToolResult(
                content=f"No pude abrir '{archivo.filename}': {type(exc).__name__}."
            )
        if extracted is None:
            return ToolResult(
                content=(
                    f"'{archivo.filename}' sí está guardado, pero su formato todavía no tiene "
                    "un lector instalado. Puedo descargarlo o convertirlo a un formato compatible."
                )
            )
        text = extracted.strip()
        truncated = len(text) > _MAX_TEXT_CHARS
        visible = text[:_MAX_TEXT_CHARS]
        return ToolResult(
            content=(visible or "El archivo no contiene texto extraíble.")
            + ("\n\n[Contenido truncado por longitud.]" if truncated else ""),
            data={
                "file_id": str(file_id),
                "filename": archivo.filename,
                "mime": archivo.mime,
                "truncated": truncated,
            },
        )


def _pdf_safe_text(value: str) -> str:
    return value.encode("latin-1", errors="replace").decode("latin-1")


def _render_text_pdf(title: str, paragraphs: list[str]) -> bytes:
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=16)
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 20)
    pdf.multi_cell(0, 10, _pdf_safe_text(title))
    pdf.ln(4)
    pdf.set_font("Helvetica", size=11)
    for paragraph in paragraphs:
        pdf.multi_cell(0, 6, _pdf_safe_text(paragraph))
        pdf.ln(3)
    return bytes(pdf.output())


def _safe_pdf_filename(raw: str) -> str:
    candidate = PurePath(raw.replace("\\", "/")).name
    candidate = "".join(char for char in candidate if 32 <= ord(char) != 127).strip(" .")
    if not candidate.lower().endswith(".pdf"):
        return "pdf-editado.pdf"
    return candidate[:250] or "pdf-editado.pdf"


class EditarPdfTool(Tool):
    name = "editar_pdf"
    description = (
        "Edita un PDF adjunto sin destruir el original y entrega un PDF nuevo descargable. "
        "Puede reconstruir su texto corregido, anexar contenido, seleccionar/eliminar páginas "
        "y rotarlas. Para corregir texto, llama primero a leer_archivo y luego pasa aquí el "
        "contenido final completo."
    )
    input_schema = {
        "type": "object",
        "properties": {
            "file_id": {"type": "string"},
            "modo": {"type": "string", "enum": ["reconstruir", "anexar", "paginas"]},
            "titulo": {"type": "string"},
            "parrafos": {"type": "array", "items": {"type": "string"}, "maxItems": 300},
            "paginas_conservar": {
                "type": "array",
                "items": {"type": "integer", "minimum": 1},
                "maxItems": 500,
            },
            "rotacion": {"type": "integer", "enum": [0, 90, 180, 270]},
            "nombre_salida": {"type": "string"},
        },
        "required": ["file_id", "modo"],
    }

    async def run(self, ctx: ToolContext, args: dict[str, Any]) -> ToolResult:
        file_id = parse_uuid(args.get("file_id"))
        if file_id is None:
            return ToolResult(content="'file_id' no es un identificador válido.")
        archivo = await _s3.descargar_archivo(ctx, file_id)
        if archivo is None:
            return ToolResult(content="No encontré ese PDF.")
        is_pdf = _normalized_mime(archivo.mime) == "application/pdf"
        if not is_pdf and not archivo.filename.lower().endswith(".pdf"):
            return ToolResult(content=f"'{archivo.filename}' no es un PDF.")

        modo = str(args.get("modo") or "").strip().lower()
        paragraphs = [str(item).strip() for item in args.get("parrafos") or [] if str(item).strip()]
        title = str(args.get("titulo") or PurePath(archivo.filename).stem).strip() or "Documento"
        try:
            if modo == "reconstruir":
                if not paragraphs:
                    return ToolResult(content="Para reconstruir el PDF necesito 'parrafos'.")
                output = _render_text_pdf(title, paragraphs)
            else:
                from pypdf import PdfReader, PdfWriter

                reader = PdfReader(io.BytesIO(archivo.contenido))
                writer = PdfWriter()
                selected = args.get("paginas_conservar") or list(range(1, len(reader.pages) + 1))
                rotation = int(args.get("rotacion") or 0)
                for page_number in selected:
                    index = int(page_number) - 1
                    if index < 0 or index >= len(reader.pages):
                        return ToolResult(content=f"La página {page_number} no existe.")
                    page = reader.pages[index]
                    if rotation:
                        page.rotate(rotation)
                    writer.add_page(page)
                if modo == "anexar":
                    if not paragraphs:
                        return ToolResult(content="Para anexar contenido necesito 'parrafos'.")
                    appendix = PdfReader(io.BytesIO(_render_text_pdf(title, paragraphs)))
                    for page in appendix.pages:
                        writer.add_page(page)
                elif modo != "paginas":
                    return ToolResult(content="'modo' debe ser reconstruir, anexar o paginas.")
                buffer = io.BytesIO()
                writer.write(buffer)
                output = buffer.getvalue()
        except Exception as exc:  # noqa: BLE001 - pypdf/fpdf exponen errores distintos
            return ToolResult(content=f"No pude editar el PDF: {type(exc).__name__}.")

        filename = _safe_pdf_filename(str(args.get("nombre_salida") or ""))
        output_id = await _s3.subir_resultado(
            ctx, filename=filename, mime="application/pdf", contenido=output
        )
        return ToolResult(
            content=f"Listo. Creé '{filename}' y conservé el original sin cambios.",
            data={
                "file_id": str(output_id),
                "filename": filename,
                "mime": "application/pdf",
                "source_file_id": str(file_id),
            },
        )
